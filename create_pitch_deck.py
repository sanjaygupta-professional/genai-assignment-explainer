#!/usr/bin/env python3.12
"""Create SamarthSchool 10-slide pitch deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colors ──────────────────────────────────────
TEAL = RGBColor(0x48, 0x72, 0x65)
TEAL_LIGHT = RGBColor(0xDD, 0xEF, 0xEF)
TEAL_DARK = RGBColor(0x35, 0x55, 0x4B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF7, 0xF6, 0xF5)
TEXT_PRIMARY = RGBColor(0x29, 0x2C, 0x33)
TEXT_SECONDARY = RGBColor(0x6B, 0x72, 0x80)
WARM_GRAY = RGBColor(0x89, 0x85, 0x84)
COPPER = RGBColor(0xBC, 0x97, 0x6A)
GOLD = RGBColor(0xD4, 0xC6, 0x8B)
RED_ACCENT = RGBColor(0xC0, 0x39, 0x2B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_bg(slide, color=WHITE):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.2):
    """Add a text box with single paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=13,
                    color=TEXT_PRIMARY, bullet_color=TEAL, spacing=6,
                    font_name="Calibri", line_spacing=1.3):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Use bullet character
        if isinstance(item, tuple):
            # (bold_part, rest)
            run1 = p.add_run()
            run1.text = item[0]
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = TEAL
            run1.font.bold = True
            run1.font.name = font_name
            run2 = p.add_run()
            run2.text = item[1]
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = font_name
        else:
            p.text = f"  {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = font_name

        p.space_after = Pt(spacing)
        p.line_spacing = Pt(font_size * line_spacing)

    return txBox


def add_slide_number(slide, num):
    """Add slide number in bottom-right."""
    add_text_box(slide, Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.35),
                 str(num), font_size=10, color=WARM_GRAY, alignment=PP_ALIGN.RIGHT)


def add_bottom_bar(slide):
    """Add thin teal bar at bottom."""
    add_rect(slide, Inches(0), Inches(7.25), SLIDE_W, Inches(0.25), TEAL)


def add_section_title(slide, title, subtitle=None):
    """Add slide title with optional subtitle."""
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
                 title, font_size=28, color=TEXT_PRIMARY, bold=True,
                 font_name="Georgia")
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.05), Inches(11), Inches(0.4),
                     subtitle, font_size=14, color=TEXT_SECONDARY)
    # Accent line under title
    add_rect(slide, Inches(0.8), Inches(1.45), Inches(1.5), Inches(0.04), TEAL)


def add_stat_card(slide, left, top, width, height, number, label, accent=TEAL):
    """Add a stat card with big number and label."""
    card = add_rect(slide, left, top, width, height, CREAM)
    # Accent bar at top of card
    add_rect(slide, left, top, width, Inches(0.05), accent)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15),
                 width - Inches(0.4), Inches(0.5),
                 number, font_size=26, color=accent, bold=True,
                 font_name="Georgia", alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), top + Inches(0.65),
                 width - Inches(0.2), Inches(0.6),
                 label, font_size=11, color=TEXT_SECONDARY,
                 alignment=PP_ALIGN.CENTER, line_spacing=1.2)


def add_table_slide(slide, left, top, width, headers, rows, col_widths=None, font_size=11):
    """Add a formatted table."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, Inches(0.4 * num_rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = "Calibri"

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else CREAM
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = TEXT_PRIMARY
                p.font.name = "Calibri"

    return table_shape


# ════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, WHITE)

# Left teal panel
add_rect(slide, Inches(0), Inches(0), Inches(5.5), SLIDE_H, TEAL)

# Title on teal panel
add_text_box(slide, Inches(0.8), Inches(1.8), Inches(4.2), Inches(1.2),
             "SamarthSchool", font_size=44, color=WHITE, bold=True,
             font_name="Georgia")
add_text_box(slide, Inches(0.8), Inches(2.9), Inches(4.2), Inches(0.5),
             "समर्थ स्कूल", font_size=24, color=RGBColor(0xDD, 0xEF, 0xEF),
             font_name="Georgia")

# Tagline on teal
add_text_box(slide, Inches(0.8), Inches(3.8), Inches(4.2), Inches(0.8),
             "AI-Powered Benefits Navigator\nfor Children with Special Abilities",
             font_size=16, color=RGBColor(0xCC, 0xDD, 0xDD), line_spacing=1.5)

# Accent line
add_rect(slide, Inches(0.8), Inches(3.5), Inches(2), Inches(0.04), GOLD)

# Right side — key stats preview
add_text_box(slide, Inches(6.5), Inches(1.5), Inches(6), Inches(0.4),
             "The Problem We Solve", font_size=20, color=TEXT_PRIMARY, bold=True,
             font_name="Georgia")

add_stat_card(slide, Inches(6.5), Inches(2.2), Inches(2.8), Inches(1.2),
              "83.6%", "of eligible families\nunaware of schemes")
add_stat_card(slide, Inches(9.7), Inches(2.2), Inches(2.8), Inches(1.2),
              "50+", "government schemes\nfragmented across portals", accent=COPPER)

add_stat_card(slide, Inches(6.5), Inches(3.7), Inches(2.8), Inches(1.2),
              "42%", "never apply because\nthey don't know", accent=RED_ACCENT)
add_stat_card(slide, Inches(9.7), Inches(3.7), Inches(2.8), Inches(1.2),
              "<5 min", "to match a child\nwith SamarthSchool", accent=TEAL)

# Footer
add_text_box(slide, Inches(6.5), Inches(6.2), Inches(6), Inches(0.4),
             "Group Assignment — Gen AI: Pre-Trained Models (Course 8919)",
             font_size=11, color=TEXT_SECONDARY)
add_text_box(slide, Inches(6.5), Inches(6.5), Inches(6), Inches(0.3),
             "GGU DBA Program via upGrad",
             font_size=10, color=WARM_GRAY)

add_bottom_bar(slide)


# ════════════════════════════════════════════════
# SLIDE 2: The Problem
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, "The Problem", "India's disability benefits: generous policy, poor awareness, low uptake")

# Three-column layout
# Column 1: The gap
add_rect(slide, Inches(0.8), Inches(1.8), Inches(3.6), Inches(5.0), CREAM)
add_text_box(slide, Inches(1.0), Inches(1.9), Inches(3.2), Inches(0.35),
             "THE AWARENESS GAP", font_size=12, color=TEAL, bold=True)

bullets1 = [
    ("83.6% ", "of eligible families unaware of scholarship schemes (IJPMR 2025)"),
    ("42% ", "never apply — they don't know benefits exist (NILERD)"),
    ("71.2% ", "lack a disability certificate — prerequisite for all schemes"),
    ("27% ", "of children with disabilities (age 5-19) never enrolled in school"),
]
add_bullet_list(slide, Inches(1.0), Inches(2.4), Inches(3.2), Inches(3.5),
                bullets1, font_size=12, spacing=10)

# Column 2: The scale
add_rect(slide, Inches(4.7), Inches(1.8), Inches(3.6), Inches(5.0), CREAM)
add_text_box(slide, Inches(4.9), Inches(1.9), Inches(3.2), Inches(0.35),
             "THE SCALE", font_size=12, color=TEAL, bold=True)

bullets2 = [
    ("63.3M ", "persons with disabilities (NFHS-5 prevalence estimate)"),
    ("18-20M ", "children with disabilities in India"),
    ("2.27M ", "CWSN enrolled in schools (UDISE+ FY2022)"),
    ("Rs 800 Cr ", "annual Samagra Shiksha CWSN allocation — much underutilized"),
]
add_bullet_list(slide, Inches(4.9), Inches(2.4), Inches(3.2), Inches(3.5),
                bullets2, font_size=12, spacing=10)

# Column 3: Why it persists
add_rect(slide, Inches(8.6), Inches(1.8), Inches(4.0), Inches(5.0), CREAM)
add_text_box(slide, Inches(8.8), Inches(1.9), Inches(3.6), Inches(0.35),
             "WHY THE GAP PERSISTS", font_size=12, color=TEAL, bold=True)

bullets3 = [
    "50+ schemes across multiple ministries",
    "Information in PDFs, circulars, gazettes",
    "21 disability categories (RPwD 2016)",
    "Complex eligibility: disability type × % × age × income × state × gender",
    "22 official languages — content mostly in English/Hindi",
    "Schemes change via budget revisions & circulars",
]
add_bullet_list(slide, Inches(8.8), Inches(2.4), Inches(3.6), Inches(3.5),
                [f"  {b}" for b in bullets3], font_size=12, spacing=8)

add_slide_number(slide, 2)
add_bottom_bar(slide)


# ════════════════════════════════════════════════
# SLIDE 3: Why Now — The Gen AI Opportunity
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, "Why Now — The Gen AI Opportunity",
                  "Existing solutions failed. Gen AI changes the equation.")

# Left: Failed approaches
add_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.45), RGBColor(0xF0, 0xE0, 0xE0))
add_text_box(slide, Inches(1.0), Inches(1.85), Inches(5.0), Inches(0.35),
             "WHAT HAS BEEN TRIED (AND FAILED)", font_size=12, color=RED_ACCENT, bold=True)

failed = [
    ("MyScheme (myscheme.gov.in): ", "Generic, form-based, 15+ fields. Solves matching but not discovery or guidance."),
    ("UDID Portal: ", "Registration system, not a navigator. Creates the ID but doesn't tell families what to do with it."),
    ("Haqdarshak (40M+ users): ", "Not disability-specialized. Lacks depth for 21 RPwD categories. Per-transaction fee model."),
    ("NGO Workshops: ", "Effective but not scalable. Reaches 50-100 families per event."),
]
add_bullet_list(slide, Inches(1.0), Inches(2.4), Inches(5.3), Inches(3.5),
                failed, font_size=12, spacing=10)

# Right: Gen AI enablers
add_rect(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(0.45), TEAL_LIGHT)
add_text_box(slide, Inches(7.0), Inches(1.85), Inches(5.4), Inches(0.35),
             "WHAT GEN AI ENABLES", font_size=12, color=TEAL, bold=True)

enables = [
    ("Natural Language Queries: ", "\"10-year-old girl, hearing impairment, Tamil Nadu government school\" → personalized results"),
    ("Knowledge Graph Reasoning: ", "Structured eligibility matching across all 21 disability categories × 50+ schemes"),
    ("Multilingual Generation: ", "bge-m3 cross-lingual embeddings + IndicTrans2 for 22 languages"),
    ("Continuous Updates: ", "Automated crawling + change detection keeps scheme data current"),
]
add_bullet_list(slide, Inches(7.0), Inches(2.4), Inches(5.4), Inches(3.5),
                enables, font_size=12, spacing=10)

# Bottom insight box
add_rect(slide, Inches(0.8), Inches(5.8), Inches(11.8), Inches(1.0), TEAL)
add_text_box(slide, Inches(1.2), Inches(5.85), Inches(11.0), Inches(0.9),
             "Key Insight: Existing solutions treat this as a SEARCH problem. SamarthSchool treats it as a\n"
             "NAVIGATION problem — \"given this child's situation, here is everything they're entitled to, how to apply,\n"
             "and what documents to prepare.\"",
             font_size=14, color=WHITE, bold=False, font_name="Calibri", line_spacing=1.4)

add_slide_number(slide, 3)
add_bottom_bar(slide)


# ════════════════════════════════════════════════
# SLIDE 4: Solution — GraphRAG Architecture
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, "Solution — GraphRAG Architecture",
                  "Structured reasoning via Knowledge Graph + natural language via RAG")

# Architecture flow — simplified visual with boxes and arrows
# Data Layer
add_rect(slide, Inches(0.8), Inches(1.9), Inches(2.5), Inches(1.8), CREAM)
add_text_box(slide, Inches(0.9), Inches(1.95), Inches(2.3), Inches(0.3),
             "DATA INGESTION", font_size=10, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)
items_data = ["Web Crawler (Scrapy)", "PDF Parser (Docling + Surya)", "Gazette RSS Monitor", "Manual Curation"]
add_bullet_list(slide, Inches(0.9), Inches(2.3), Inches(2.3), Inches(1.2),
                [f"  {x}" for x in items_data], font_size=10, spacing=2)

# Arrow
add_text_box(slide, Inches(3.4), Inches(2.5), Inches(0.5), Inches(0.4),
             "→", font_size=28, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

# Storage Layer
add_rect(slide, Inches(3.9), Inches(1.9), Inches(2.8), Inches(1.8), TEAL_LIGHT)
add_text_box(slide, Inches(4.0), Inches(1.95), Inches(2.6), Inches(0.3),
             "STORAGE", font_size=10, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)
items_store = [
    ("Neo4j: ", "Knowledge Graph\n    (scheme eligibility rules)"),
    ("Qdrant: ", "Vector Store\n    (document embeddings, bge-m3)"),
    ("PostgreSQL: ", "User data, audit logs"),
]
add_bullet_list(slide, Inches(4.0), Inches(2.3), Inches(2.6), Inches(1.2),
                items_store, font_size=10, spacing=4)

# Arrow
add_text_box(slide, Inches(6.8), Inches(2.5), Inches(0.5), Inches(0.4),
             "→", font_size=28, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

# Query Layer
add_rect(slide, Inches(7.3), Inches(1.9), Inches(2.8), Inches(1.8), CREAM)
add_text_box(slide, Inches(7.4), Inches(1.95), Inches(2.6), Inches(0.3),
             "QUERY ROUTING", font_size=10, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)
items_query = [
    ("Eligibility: ", "→ KG (Cypher query)"),
    ("Explanation: ", "→ RAG (vector search)"),
    ("Complex: ", "→ Hybrid (both)"),
]
add_bullet_list(slide, Inches(7.4), Inches(2.3), Inches(2.6), Inches(1.2),
                items_query, font_size=10, spacing=4)

# Arrow
add_text_box(slide, Inches(10.2), Inches(2.5), Inches(0.5), Inches(0.4),
             "→", font_size=28, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

# Response
add_rect(slide, Inches(10.7), Inches(1.9), Inches(2.0), Inches(1.8), TEAL)
add_text_box(slide, Inches(10.8), Inches(1.95), Inches(1.8), Inches(0.3),
             "RESPONSE", font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
items_resp = ["Gemini 2.0 Flash", "User's language", "Citations", "Confidence score", "Advisory disclaimer"]
add_bullet_list(slide, Inches(10.8), Inches(2.3), Inches(1.8), Inches(1.2),
                [f"  {x}" for x in items_resp], font_size=9, spacing=2, color=WHITE)

# Tech stack table below
add_text_box(slide, Inches(0.8), Inches(4.1), Inches(5), Inches(0.35),
             "Technology Stack", font_size=16, color=TEXT_PRIMARY, bold=True, font_name="Georgia")

headers = ["Component", "Technology", "Why"]
rows = [
    ["Embedding", "BAAI/bge-m3 (1024d)", "Multilingual; dense+sparse hybrid"],
    ["Vector DB", "Qdrant (self-hosted)", "Sparse vector support; 4GB RAM"],
    ["Graph DB", "Neo4j Community 5.x", "Cypher queries; scheme ontology"],
    ["LLM", "Gemini 2.0 Flash (paid)", "Best Hindi generation; $0.075/1M tokens"],
    ["Framework", "LlamaIndex PropertyGraphIndex", "Native GraphRAG; Qdrant+Neo4j integration"],
    ["Multilingual", "IndicTrans2 + IndicXlit", "22 Indian languages; transliteration"],
    ["Frontend", "React + Next.js", "WCAG 2.1 AA; mobile-first"],
]
add_table_slide(slide, Inches(0.8), Inches(4.5), Inches(11.8), headers, rows,
                col_widths=[Inches(2), Inches(3.5), Inches(6.3)], font_size=10)

add_slide_number(slide, 4)
add_bottom_bar(slide)


# ════════════════════════════════════════════════
# SLIDE 5: How It Works
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, "How It Works", "From natural language query to personalized action plan in <5 minutes")

# Example query box
add_rect(slide, Inches(0.8), Inches(1.8), Inches(11.8), Inches(1.0), CREAM)
add_text_box(slide, Inches(1.0), Inches(1.85), Inches(1.8), Inches(0.3),
             "EXAMPLE QUERY:", font_size=11, color=TEAL, bold=True)
add_text_box(slide, Inches(1.0), Inches(2.15), Inches(11.4), Inches(0.5),
             "\"I have a 12-year-old boy with 50% locomotor disability in a government school in Karnataka. "
             "Family income is Rs 1.5 lakh per year. What benefits is he eligible for?\"",
             font_size=13, color=TEXT_PRIMARY, font_name="Georgia")

# Flow steps
steps = [
    ("1", "DESCRIBE", "School admin describes\nchild's situation in\nnatural language\n(any Indian language)"),
    ("2", "DETECT", "System detects language\n(fastText lid.176)\nTranslates for KG query\nKeeps original for RAG"),
    ("3", "MATCH", "Knowledge Graph filters\neligible schemes via\nCypher: disability type ×\n% × age × income × state"),
    ("4", "EXPLAIN", "RAG retrieves scheme\ndocuments, generates\nexplanations, documents\nneeded, application steps"),
    ("5", "RESPOND", "Complete action plan\nin user's language:\n8-12 matching schemes\nwith citations & disclaimer"),
]

for i, (num, title, desc) in enumerate(steps):
    left = Inches(0.8 + i * 2.5)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.8), Inches(3.2),
                                     Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = TEAL
    circle.line.fill.background()
    # Number text
    add_text_box(slide, left + Inches(0.85), Inches(3.22), Inches(0.4), Inches(0.4),
                 num, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, left, Inches(3.8), Inches(2.2), Inches(0.35),
                 title, font_size=13, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    add_text_box(slide, left, Inches(4.15), Inches(2.2), Inches(1.5),
                 desc, font_size=11, color=TEXT_SECONDARY, alignment=PP_ALIGN.CENTER,
                 line_spacing=1.3)

    # Arrow between steps
    if i < 4:
        add_text_box(slide, left + Inches(2.1), Inches(3.3), Inches(0.4), Inches(0.4),
                     "→", font_size=20, color=WARM_GRAY, alignment=PP_ALIGN.CENTER)

# Result box
add_rect(slide, Inches(0.8), Inches(5.8), Inches(11.8), Inches(1.1), TEAL_LIGHT)
add_text_box(slide, Inches(1.0), Inches(5.85), Inches(11.4), Inches(0.3),
             "RESULT FOR THE EXAMPLE QUERY:", font_size=11, color=TEAL, bold=True)
result_items = (
    "Pre-Matric Scholarship (Rs 6,000 + Rs 1,000 books)  •  Samagra Shiksha CWSN (Rs 3,500/yr)  •  "
    "ADIP Assistive Device (Rs 15,000)  •  Karnataka State Disability Pension  •  Niramaya Health Insurance  •  "
    "National Trust Schemes  •  ... + 2-6 more state-specific schemes — with documents needed, portal links, "
    "and step-by-step application guidance"
)
add_text_box(slide, Inches(1.0), Inches(6.15), Inches(11.4), Inches(0.6),
             result_items, font_size=11, color=TEXT_PRIMARY, line_spacing=1.3)

add_slide_number(slide, 5)
add_bottom_bar(slide)


# ════════════════════════════════════════════════
# SLIDE 6: Product Roadmap
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, "Product Roadmap", "From MVP to national scale in 36 months")

# Timeline phases as horizontal cards
phases = [
    ("PHASE 1: MVP", "Months 0-3", [
        "50+ central schemes in KG",
        "English + Hindi",
        "Basic web interface",
        "RAGAS evaluation suite",
        "5 pilot schools in 1 city",
    ], "5\nschools"),
    ("PHASE 2: BETA", "Months 3-6", [
        "+5 states (MH, KA, TN, DL, RJ)",
        "+4 languages (Tamil, Kannada, Marathi)",
        "Document checklist generator",
        "Application guidance module",
        "25 schools across 3 cities",
    ], "25\nschools"),
    ("PHASE 3: v1.0", "Months 6-12", [
        "10 states (80%+ PwD coverage)",
        "DPDP Act + GIGW compliance",
        "WhatsApp chatbot interface",
        "Analytics dashboard",
        "Offline district-hub mode",
        "100 schools, 5 states",
    ], "100\nschools"),
    ("PHASE 4: SCALE", "Years 1-3", [
        "All 28 states + 8 UTs",
        "B2G integration (UDID, Samagra)",
        "Agentic workflows (advisory)",
        "KG API for researchers/NGOs",
        "Bhashini 22-language support",
        "10,000+ schools",
    ], "10,000+\nschools"),
]

for i, (title, period, items, metric) in enumerate(phases):
    left = Inches(0.5 + i * 3.1)
    w = Inches(2.9)

    # Phase card
    add_rect(slide, left, Inches(1.8), w, Inches(4.8), CREAM)
    # Header
    add_rect(slide, left, Inches(1.8), w, Inches(0.7), TEAL if i == 0 else TEAL_DARK)
    add_text_box(slide, left + Inches(0.1), Inches(1.82), w - Inches(0.2), Inches(0.3),
                 title, font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(2.1), w - Inches(0.2), Inches(0.3),
                 period, font_size=10, color=RGBColor(0xCC, 0xDD, 0xDD), alignment=PP_ALIGN.CENTER)

    # Items
    add_bullet_list(slide, left + Inches(0.15), Inches(2.65), w - Inches(0.3), Inches(3.0),
                    [f"  {x}" for x in items], font_size=10, spacing=4, line_spacing=1.25)

    # Metric badge
    add_rect(slide, left + Inches(0.6), Inches(5.9), Inches(1.7), Inches(0.6), TEAL_LIGHT)
    add_text_box(slide, left + Inches(0.6), Inches(5.92), Inches(1.7), Inches(0.55),
                 metric, font_size=12, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

    # Arrow between phases
    if i < 3:
        add_text_box(slide, left + w - Inches(0.1), Inches(3.5), Inches(0.4), Inches(0.4),
                     "→", font_size=22, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 6)
add_bottom_bar(slide)


# ════════════════════════════════════════════════
# SLIDE 7: Adoption Strategy
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, "Adoption Strategy",
                  "Go through existing networks, not cold outreach")

# 5 strategy cards in 2 rows
strategies = [
    ("CWSN Coordinators", "Channel",
     "Every district has a CWSN coordinator under Samagra Shiksha whose job is ensuring benefit uptake. "
     "SamarthSchool becomes their tool — aligned incentives."),
    ("NGO Amplification", "Partners",
     "Samarthanam Trust, CBM India, Sense International already do school outreach. "
     "SamarthSchool is a tool they use during visits, not a standalone product."),
    ("WhatsApp-First", "Interface",
     "500M+ Indian WhatsApp users. No app install needed. "
     "A CWSN coordinator queries for multiple schools from their phone."),
]

strategies2 = [
    ("Pilot Evidence", "Sales Collateral",
     "5-school MVP generates before/after data: \"20 CWSN accessed 6 benefits → now 38.\" "
     "This evidence sells the product."),
    ("Navigator vs. Search", "vs. MyScheme",
     "MyScheme: 15+ form fields, generic, user must know to search. "
     "SamarthSchool: describe the child, get a complete action plan. UX gap = adoption driver."),
]

for i, (title, tag, desc) in enumerate(strategies):
    left = Inches(0.8 + i * 4.0)
    add_rect(slide, left, Inches(1.8), Inches(3.7), Inches(2.3), CREAM)
    # Tag
    add_rect(slide, left + Inches(0.15), Inches(1.95), Inches(1.2), Inches(0.3), TEAL)
    add_text_box(slide, left + Inches(0.15), Inches(1.95), Inches(1.2), Inches(0.28),
                 tag, font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, left + Inches(0.15), Inches(2.35), Inches(3.4), Inches(0.35),
                 title, font_size=15, color=TEXT_PRIMARY, bold=True, font_name="Georgia")
    # Description
    add_text_box(slide, left + Inches(0.15), Inches(2.7), Inches(3.4), Inches(1.2),
                 desc, font_size=11, color=TEXT_SECONDARY, line_spacing=1.3)

for i, (title, tag, desc) in enumerate(strategies2):
    left = Inches(0.8 + i * 6.0)
    w = Inches(5.7) if i == 0 else Inches(5.7)
    add_rect(slide, left, Inches(4.4), Inches(5.7), Inches(2.2), CREAM)
    # Tag
    add_rect(slide, left + Inches(0.15), Inches(4.55), Inches(1.8), Inches(0.3), COPPER)
    add_text_box(slide, left + Inches(0.15), Inches(4.55), Inches(1.8), Inches(0.28),
                 tag, font_size=9, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_text_box(slide, left + Inches(0.15), Inches(4.95), Inches(5.4), Inches(0.35),
                 title, font_size=15, color=TEXT_PRIMARY, bold=True, font_name="Georgia")
    # Description
    add_text_box(slide, left + Inches(0.15), Inches(5.35), Inches(5.4), Inches(1.0),
                 desc, font_size=11, color=TEXT_SECONDARY, line_spacing=1.3)

add_slide_number(slide, 7)
add_bottom_bar(slide)


# ════════════════════════════════════════════════
# SLIDE 8: ROI & Impact
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, "ROI & Impact",
                  "Conservative projections with derived benefit values")

# Impact table
headers = ["Metric", "Year 1", "Year 2", "Year 3"]
rows = [
    ["Schools served", "50", "200", "500"],
    ["Children impacted", "750", "3,000", "7,500"],
    ["Additional benefits/child", "1.7", "2.0", "2.5"],
    ["Total benefit applications", "1,275", "6,000", "18,750"],
    ["Value of benefits applied for", "INR 41 L", "INR 1.92 Cr", "INR 6.0 Cr"],
    ["Total cost", "INR 2.09 Cr", "INR 2.46 Cr", "INR 2.75 Cr"],
    ["Social ROI", "0.20x", "0.78x", "2.2x"],
]
add_table_slide(slide, Inches(0.8), Inches(1.7), Inches(7.5), headers, rows,
                col_widths=[Inches(2.8), Inches(1.5), Inches(1.5), Inches(1.7)], font_size=11)

# Right side: Key assumptions
add_text_box(slide, Inches(8.8), Inches(1.7), Inches(4), Inches(0.35),
             "Key Assumptions", font_size=14, color=TEXT_PRIMARY, bold=True, font_name="Georgia")

assumptions = [
    ("Rs 3,200 ", "weighted avg benefit value\n(derived from actual scheme amounts)"),
    ("15 CWSN ", "per school average\n(from UDISE+ data)"),
    ("50→200→500 ", "school growth\n(conservative; B2G takes 12-24 months)"),
]
add_bullet_list(slide, Inches(8.8), Inches(2.2), Inches(4.0), Inches(2.5),
                assumptions, font_size=11, spacing=12)

# Break-even box
add_rect(slide, Inches(8.8), Inches(4.3), Inches(4.0), Inches(1.5), CREAM)
add_text_box(slide, Inches(9.0), Inches(4.4), Inches(3.6), Inches(0.3),
             "BREAK-EVEN", font_size=12, color=TEAL, bold=True)
add_text_box(slide, Inches(9.0), Inches(4.75), Inches(3.6), Inches(0.9),
             "Year 4-5 with B2G contracts from 2-3 state education departments. "
             "Consistent with social venture norms — Haqdarshak operated on grant/impact funding for 5+ years.",
             font_size=11, color=TEXT_SECONDARY, line_spacing=1.3)

# Sensitivity box
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(5), Inches(0.35),
             "Sensitivity Analysis", font_size=14, color=TEXT_PRIMARY, bold=True, font_name="Georgia")

sens_headers = ["Scenario", "Schools (Y3)", "Revenue (Y3)", "Net Position"]
sens_rows = [
    ["Pessimistic (50% adoption)", "250", "INR 1.0 Cr", "-INR 1.75 Cr"],
    ["Base case", "500", "INR 2.0 Cr", "-INR 0.75 Cr"],
    ["Optimistic (B2G in Y2)", "1,000", "INR 4.0 Cr", "+INR 1.25 Cr"],
]
add_table_slide(slide, Inches(0.8), Inches(5.9), Inches(7.5), sens_headers, sens_rows,
                col_widths=[Inches(2.8), Inches(1.5), Inches(1.5), Inches(1.7)], font_size=10)

add_slide_number(slide, 8)
add_bottom_bar(slide)


# ════════════════════════════════════════════════
# SLIDE 9: Competitive Moat & Unique Aspects
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, "Competitive Moat & Unique Aspects")

# Three columns
# Col 1: Data Moat
add_rect(slide, Inches(0.8), Inches(1.8), Inches(3.7), Inches(3.5), CREAM)
add_rect(slide, Inches(0.8), Inches(1.8), Inches(3.7), Inches(0.05), TEAL)
add_text_box(slide, Inches(1.0), Inches(2.0), Inches(3.3), Inches(0.35),
             "DATA MOAT", font_size=14, color=TEAL, bold=True, font_name="Georgia")
moat_items = [
    "First machine-queryable ontology of Indian disability welfare schemes",
    "50+ schemes with eligibility rules as graph-traversable entities",
    "Deepens over time: state expansion, validation, usage patterns",
    "Requires domain expertise + NGO partnerships + continuous curation",
    "Not easily replicable by competitors",
]
add_bullet_list(slide, Inches(1.0), Inches(2.45), Inches(3.3), Inches(2.5),
                [f"  {x}" for x in moat_items], font_size=11, spacing=6)

# Col 2: IP & Research
add_rect(slide, Inches(4.8), Inches(1.8), Inches(3.7), Inches(3.5), CREAM)
add_rect(slide, Inches(4.8), Inches(1.8), Inches(3.7), Inches(0.05), COPPER)
add_text_box(slide, Inches(5.0), Inches(2.0), Inches(3.3), Inches(0.35),
             "IP & RESEARCH", font_size=14, color=COPPER, bold=True, font_name="Georgia")
ip_items = [
    "Patent potential: structured eligibility determination for Indian welfare",
    "Novel domain application (not novel technique)",
    "Publishable at: ICTD, ACM DEV, AAAI AI4SG Workshop, CHI LBW",
    "Benchmark dataset: 200+ eligibility queries with ground-truth",
    "GraphRAG vs. pure RAG evaluation for policy matching",
]
add_bullet_list(slide, Inches(5.0), Inches(2.45), Inches(3.3), Inches(2.5),
                [f"  {x}" for x in ip_items], font_size=11, spacing=6)

# Col 3: Market Position
add_rect(slide, Inches(8.8), Inches(1.8), Inches(3.7), Inches(3.5), CREAM)
add_rect(slide, Inches(8.8), Inches(1.8), Inches(3.7), Inches(0.05), GOLD)
add_text_box(slide, Inches(9.0), Inches(2.0), Inches(3.3), Inches(0.35),
             "MARKET POSITION", font_size=14, color=RGBColor(0xA0, 0x90, 0x50), bold=True, font_name="Georgia")
market_items = [
    "Intersection: Disability inclusion × GovTech × EdTech",
    "CSR pool: Rs 29,987 Cr ($3.6B), 44% to education/disability",
    "Impact investors: AssisTech Foundation, Aavishkaar, Omidyar",
    "Social Stock Exchange listing pathway",
    "India EdTech: $3.6-12.1B at 27% CAGR",
]
add_bullet_list(slide, Inches(9.0), Inches(2.45), Inches(3.3), Inches(2.5),
                [f"  {x}" for x in market_items], font_size=11, spacing=6)

# Comparison table
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(5), Inches(0.35),
             "Competitive Comparison", font_size=14, color=TEXT_PRIMARY, bold=True, font_name="Georgia")

comp_headers = ["Feature", "MyScheme", "Haqdarshak", "SamarthSchool"]
comp_rows = [
    ["Disability-specialized", "No", "No", "Yes (21 categories)"],
    ["Natural language queries", "No (form-based)", "Limited", "Yes (multilingual)"],
    ["Knowledge Graph reasoning", "Rule-based", "No", "Yes (Neo4j + Cypher)"],
    ["School-admin workflow", "No", "No", "Yes (designed for schools)"],
    ["Eligibility + guidance", "Eligibility only", "Partial", "Full action plan"],
]
add_table_slide(slide, Inches(0.8), Inches(5.9), Inches(11.8), comp_headers, comp_rows,
                col_widths=[Inches(3), Inches(2.5), Inches(2.5), Inches(3.8)], font_size=10)

add_slide_number(slide, 9)
add_bottom_bar(slide)


# ════════════════════════════════════════════════
# SLIDE 10: The Ask
# ════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# Left panel — teal
add_rect(slide, Inches(0), Inches(0), Inches(5.5), SLIDE_H, TEAL)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(4.2), Inches(0.6),
             "The Ask", font_size=36, color=WHITE, bold=True, font_name="Georgia")
add_rect(slide, Inches(0.8), Inches(1.9), Inches(2), Inches(0.04), GOLD)

# Team
add_text_box(slide, Inches(0.8), Inches(2.3), Inches(4.2), Inches(0.3),
             "TEAM", font_size=11, color=RGBColor(0xCC, 0xDD, 0xDD), bold=True)
team_items = [
    "2 ML Engineers (RAG/KG)",
    "1 Full-Stack Developer",
    "1 Disability Domain Expert",
    "1 Policy Researcher / KG Curator",
    "1 Product Manager",
    "0.5 DevOps (part-time)",
]
add_bullet_list(slide, Inches(0.8), Inches(2.65), Inches(4.2), Inches(2.2),
                [f"  {x}" for x in team_items], font_size=12, spacing=4, color=WHITE)

# Funding
add_text_box(slide, Inches(0.8), Inches(4.6), Inches(4.2), Inches(0.3),
             "FUNDING", font_size=11, color=RGBColor(0xCC, 0xDD, 0xDD), bold=True)
add_text_box(slide, Inches(0.8), Inches(4.95), Inches(4.2), Inches(0.6),
             "INR 2.1 Crore (Year 1)\nCSR/grant funded — Tata Trusts, Infosys Foundation, Wipro",
             font_size=13, color=WHITE, line_spacing=1.4)

# Timeline
add_text_box(slide, Inches(0.8), Inches(5.7), Inches(4.2), Inches(0.3),
             "TIMELINE", font_size=11, color=RGBColor(0xCC, 0xDD, 0xDD), bold=True)
add_text_box(slide, Inches(0.8), Inches(6.05), Inches(4.2), Inches(0.6),
             "MVP in 3 months → Pilot evidence in 6 months\n→ 100 schools by Month 12",
             font_size=13, color=WHITE, line_spacing=1.4)

# Right panel — the close
add_text_box(slide, Inches(6.2), Inches(1.5), Inches(6.5), Inches(0.5),
             "The Success Metric", font_size=22, color=TEXT_PRIMARY, bold=True, font_name="Georgia")

add_rect(slide, Inches(6.2), Inches(2.2), Inches(6.5), Inches(2.0), CREAM)
add_text_box(slide, Inches(6.5), Inches(2.4), Inches(6.0), Inches(1.6),
             "If 5 pilot schools demonstrate a measurable increase in benefits accessed by children "
             "with special abilities, the evidence unlocks everything:\n\n"
             "→ B2G procurement from state education departments\n"
             "→ Impact investment (AssisTech, Aavishkaar, Omidyar)\n"
             "→ National scale via Samagra Shiksha network",
             font_size=13, color=TEXT_PRIMARY, line_spacing=1.4)

# Closing quote
add_rect(slide, Inches(6.2), Inches(4.8), Inches(6.5), Inches(1.8), TEAL)
add_text_box(slide, Inches(6.5), Inches(5.0), Inches(6.0), Inches(1.4),
             "\"Every child with special abilities\ndeserves every benefit they're entitled to.\n\n"
             "SamarthSchool makes that possible.\"",
             font_size=18, color=WHITE, font_name="Georgia",
             alignment=PP_ALIGN.CENTER, line_spacing=1.4)

add_slide_number(slide, 10)
add_bottom_bar(slide)


# ── Save ────────────────────────────────────────
output_path = "/home/opc/genai-assignment-explainer/SamarthSchool-PitchDeck.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
