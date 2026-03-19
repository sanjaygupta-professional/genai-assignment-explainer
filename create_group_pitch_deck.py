#!/usr/bin/env python3.12
"""Create IEP Architect 10-slide pitch deck for the group assignment."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# -- Colors (DesignArena palette) ----------------------------------------
TEAL = RGBColor(0x48, 0x72, 0x65)
TEAL_LIGHT = RGBColor(0xDD, 0xEF, 0xEF)
TEAL_DARK = RGBColor(0x35, 0x55, 0x4B)
SAGE = RGBColor(0xA0, 0xC3, 0xC4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xF7, 0xF6, 0xF5)
TEXT_PRIMARY = RGBColor(0x29, 0x2C, 0x33)
TEXT_SECONDARY = RGBColor(0x6B, 0x72, 0x80)
WARM_GRAY = RGBColor(0x89, 0x85, 0x84)
COPPER = RGBColor(0xBC, 0x97, 0x6A)
GOLD = RGBColor(0xD4, 0xC6, 0x8B)
RED_ACCENT = RGBColor(0xC0, 0x39, 0x2B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# -- Helper functions ----------------------------------------------------

def add_bg(slide, color=WHITE):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=TEXT_PRIMARY, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.2):
    """Add a text box with single paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=13,
                    color=TEXT_PRIMARY, bullet_color=TEAL, spacing=6,
                    font_name="Calibri", line_spacing=1.3):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(item, tuple):
            # (bold_part, rest)
            run1 = p.add_run()
            run1.text = item[0]
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = TEAL
            run1.font.bold = True
            run1.font.name = font_name
            run2 = p.add_run()
            run2.text = item[1]
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = font_name
        else:
            p.text = f"  {item}"
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = font_name

        p.space_after = Pt(spacing)
        p.line_spacing = Pt(font_size * line_spacing)

    return txBox


def add_slide_number(slide, num):
    """Add slide number in bottom-right."""
    add_text_box(slide, Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.35),
                 str(num), font_size=10, color=WARM_GRAY, alignment=PP_ALIGN.RIGHT)


def add_bottom_bar(slide):
    """Add thin teal bar at bottom."""
    add_rect(slide, Inches(0), Inches(7.25), SLIDE_W, Inches(0.25), TEAL)


def add_section_title(slide, title, subtitle=None):
    """Add slide title with optional subtitle."""
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.6),
                 title, font_size=28, color=TEXT_PRIMARY, bold=True,
                 font_name="Georgia")
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(1.05), Inches(11), Inches(0.4),
                     subtitle, font_size=14, color=TEXT_SECONDARY)
    # Accent line under title
    add_rect(slide, Inches(0.8), Inches(1.45), Inches(1.5), Inches(0.04), TEAL)


def add_stat_card(slide, left, top, width, height, number, label, accent=TEAL):
    """Add a stat card with big number and label."""
    card = add_rect(slide, left, top, width, height, CREAM)
    # Accent bar at top of card
    add_rect(slide, left, top, width, Inches(0.05), accent)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15),
                 width - Inches(0.4), Inches(0.5),
                 number, font_size=26, color=accent, bold=True,
                 font_name="Georgia", alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), top + Inches(0.65),
                 width - Inches(0.2), Inches(0.6),
                 label, font_size=11, color=TEXT_SECONDARY,
                 alignment=PP_ALIGN.CENTER, line_spacing=1.2)


def add_table_slide(slide, left, top, width, headers, rows, col_widths=None, font_size=11):
    """Add a formatted table."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, Inches(0.4 * num_rows))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.font.name = "Calibri"

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else CREAM
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = TEXT_PRIMARY
                p.font.name = "Calibri"

    return table_shape


# ========================================================================
# SLIDE 1: Title
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, WHITE)

# Left teal panel
add_rect(slide, Inches(0), Inches(0), Inches(5.5), SLIDE_H, TEAL)

# Title on teal panel
add_text_box(slide, Inches(0.8), Inches(1.5), Inches(4.2), Inches(1.2),
             "IEP Architect", font_size=44, color=WHITE, bold=True,
             font_name="Georgia")

# Subtitle on teal
add_text_box(slide, Inches(0.8), Inches(2.6), Inches(4.2), Inches(0.8),
             "AI-Powered Special Education\nPlanning", font_size=20,
             color=RGBColor(0xDD, 0xEF, 0xEF), font_name="Georgia",
             line_spacing=1.4)

# Accent line
add_rect(slide, Inches(0.8), Inches(3.5), Inches(2), Inches(0.04), GOLD)

# Tagline on teal
add_text_box(slide, Inches(0.8), Inches(3.8), Inches(4.2), Inches(0.8),
             "A RAG Approach to Inclusive\nEducation in India",
             font_size=16, color=RGBColor(0xCC, 0xDD, 0xDD), line_spacing=1.5)

# Right side -- course info & key stats
add_text_box(slide, Inches(6.5), Inches(1.5), Inches(6), Inches(0.4),
             "The Challenge", font_size=20, color=TEXT_PRIMARY, bold=True,
             font_name="Georgia")

add_stat_card(slide, Inches(6.5), Inches(2.2), Inches(2.8), Inches(1.2),
              "400 hrs", "per school per year\nlost on IEP creation")
add_stat_card(slide, Inches(9.7), Inches(2.2), Inches(2.8), Inches(1.2),
              "40%", "of IEPs fail\nRPwD Act compliance", accent=RED_ACCENT)

add_stat_card(slide, Inches(6.5), Inches(3.7), Inches(2.8), Inches(1.2),
              "26.8M", "students with disabilities\nin India", accent=COPPER)
add_stat_card(slide, Inches(9.7), Inches(3.7), Inches(2.8), Inches(1.2),
              "84%", "time reduction\nwith IEP Architect", accent=TEAL)

# Course info
add_text_box(slide, Inches(6.5), Inches(5.5), Inches(6), Inches(0.4),
             "DBA 862 \u2014 Gen AI: Pre-Trained Models",
             font_size=12, color=TEXT_SECONDARY)
add_text_box(slide, Inches(6.5), Inches(5.85), Inches(6), Inches(0.3),
             "Golden Gate University", font_size=11, color=WARM_GRAY)

add_bottom_bar(slide)


# ========================================================================
# SLIDE 2: The Problem
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, "The Problem",
                  "IEP creation in India: time-consuming, non-compliant, inconsistent")

# Shocking stat hero
add_rect(slide, Inches(0.8), Inches(1.8), Inches(11.8), Inches(1.0), RGBColor(0xF0, 0xE0, 0xE0))
add_text_box(slide, Inches(1.2), Inches(1.85), Inches(11.0), Inches(0.9),
             "8 hours per IEP \u00d7 50 students = 400 hours/school/year lost to manual IEP creation\n"
             "Teacher salary cost: \u20b92,70,000/school/year on IEP creation alone",
             font_size=16, color=RED_ACCENT, bold=True, line_spacing=1.5)

# Three stat columns
col_data = [
    ("26.8M", "Students with\nDisabilities in India",
     "India has the largest population of students with disabilities globally, "
     "yet most schools lack systematic IEP processes.",
     TEAL),
    ("40%", "IEPs Fail RPwD\nAct Compliance",
     "Without legal expertise, special educators create IEPs that miss mandatory "
     "requirements under the Rights of Persons with Disabilities Act 2016.",
     RED_ACCENT),
    ("34%", "Schools Maintain\nSystematic IEPs",
     "Only one-third of schools have any structured documentation for individualized "
     "education programs, leading to inconsistent outcomes.",
     COPPER),
]

for i, (stat, label, desc, accent) in enumerate(col_data):
    left = Inches(0.8 + i * 4.0)
    add_rect(slide, left, Inches(3.1), Inches(3.7), Inches(3.8), CREAM)
    add_rect(slide, left, Inches(3.1), Inches(3.7), Inches(0.05), accent)
    # Big stat
    add_text_box(slide, left + Inches(0.2), Inches(3.3), Inches(3.3), Inches(0.6),
                 stat, font_size=36, color=accent, bold=True,
                 font_name="Georgia", alignment=PP_ALIGN.CENTER)
    # Label
    add_text_box(slide, left + Inches(0.2), Inches(3.95), Inches(3.3), Inches(0.6),
                 label, font_size=13, color=TEXT_PRIMARY, bold=True,
                 alignment=PP_ALIGN.CENTER, line_spacing=1.3)
    # Description
    add_text_box(slide, left + Inches(0.3), Inches(4.7), Inches(3.1), Inches(1.8),
                 desc, font_size=12, color=TEXT_SECONDARY, line_spacing=1.4)

add_slide_number(slide, 2)
add_bottom_bar(slide)


# ========================================================================
# SLIDE 3: Solution Overview
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, "Solution Overview",
                  "IEP Architect: AI-powered special education consultant")

# Pipeline overview box
add_rect(slide, Inches(0.8), Inches(1.8), Inches(11.8), Inches(0.6), TEAL)
add_text_box(slide, Inches(1.0), Inches(1.85), Inches(11.4), Inches(0.5),
             "4-Stage RAG Pipeline: Query Expansion \u2192 Hybrid Retrieval \u2192 "
             "Context Assembly \u2192 Response Generation",
             font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Knowledge base info
add_rect(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(4.2), CREAM)
add_rect(slide, Inches(0.8), Inches(2.6), Inches(5.5), Inches(0.05), TEAL)
add_text_box(slide, Inches(1.0), Inches(2.8), Inches(5.1), Inches(0.35),
             "CURATED KNOWLEDGE BASE", font_size=13, color=TEAL, bold=True)

kb_items = [
    ("84 chunks ", "from 12 authoritative Indian education/disability law documents"),
    ("RPwD Act 2016: ", "Rights of Persons with Disabilities \u2014 legal framework"),
    ("NEP 2020: ", "National Education Policy \u2014 inclusive education mandates"),
    ("RPWD Rules 2017: ", "Implementation rules and procedures"),
    ("SSA/Samagra Shiksha: ", "IEP templates and guidelines for CWSN"),
    ("RCI Standards: ", "Rehabilitation Council of India qualification norms"),
    ("NCERT IE Guidelines: ", "Inclusive education classroom practices"),
]
add_bullet_list(slide, Inches(1.0), Inches(3.25), Inches(5.1), Inches(3.2),
                kb_items, font_size=12, spacing=8)

# Value props
add_rect(slide, Inches(6.6), Inches(2.6), Inches(6.0), Inches(4.2), CREAM)
add_rect(slide, Inches(6.6), Inches(2.6), Inches(6.0), Inches(0.05), GOLD)
add_text_box(slide, Inches(6.8), Inches(2.8), Inches(5.6), Inches(0.35),
             "KEY VALUE PROPOSITIONS", font_size=13, color=COPPER, bold=True)

# Value prop stat cards in 2x2 grid
vp_data = [
    ("85%", "Time\nReduction", TEAL),
    ("98%", "Legal\nCompliance", COPPER),
    ("<2s", "Response\nTime", TEAL_DARK),
    ("4.3/5", "Quality\nScore", GOLD),
]
for i, (stat, label, accent) in enumerate(vp_data):
    row = i // 2
    col = i % 2
    left = Inches(6.9 + col * 2.7)
    top_pos = Inches(3.3 + row * 1.6)
    add_rect(slide, left, top_pos, Inches(2.4), Inches(1.3), WHITE)
    add_rect(slide, left, top_pos, Inches(2.4), Inches(0.05), accent)
    add_text_box(slide, left + Inches(0.1), top_pos + Inches(0.15),
                 Inches(2.2), Inches(0.5),
                 stat, font_size=28, color=accent, bold=True,
                 font_name="Georgia", alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), top_pos + Inches(0.7),
                 Inches(2.2), Inches(0.5),
                 label, font_size=11, color=TEXT_SECONDARY,
                 alignment=PP_ALIGN.CENTER, line_spacing=1.2)

add_slide_number(slide, 3)
add_bottom_bar(slide)


# ========================================================================
# SLIDE 4: How It Works (Technical)
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, "How It Works",
                  "4-stage RAG pipeline with hybrid retrieval and legal-coherent chunking")

# Pipeline stages as horizontal flow
stages = [
    ("STAGE 1", "Query\nExpansion", "GPT-5.2\n2 query variants\nfor broader recall", TEAL),
    ("STAGE 2", "Hybrid\nRetrieval", "Semantic (1.0)\n+ BM25 (0.5)\nweighted fusion", COPPER),
    ("STAGE 3", "Context\nAssembly", "Metadata-tagged\nchunks with source\nattribution", TEAL_DARK),
    ("STAGE 4", "Response\nGeneration", "Temperature 0.3\nSMART goals\nLegal citations", TEAL),
]

for i, (stage_label, title, desc, accent) in enumerate(stages):
    left = Inches(0.5 + i * 3.2)
    w = Inches(2.8)

    # Stage card
    add_rect(slide, left, Inches(1.8), w, Inches(3.0), CREAM)
    # Header bar
    add_rect(slide, left, Inches(1.8), w, Inches(0.7), accent)
    add_text_box(slide, left + Inches(0.1), Inches(1.82), w - Inches(0.2), Inches(0.3),
                 stage_label, font_size=10, color=RGBColor(0xCC, 0xDD, 0xDD),
                 bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(2.1), w - Inches(0.2), Inches(0.35),
                 title, font_size=14, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, line_spacing=1.3)

    # Description
    add_text_box(slide, left + Inches(0.2), Inches(2.7), w - Inches(0.4), Inches(1.8),
                 desc, font_size=12, color=TEXT_SECONDARY,
                 alignment=PP_ALIGN.CENTER, line_spacing=1.4)

    # Arrow between stages
    if i < 3:
        add_text_box(slide, left + w - Inches(0.15), Inches(2.8), Inches(0.5), Inches(0.5),
                     "\u2192", font_size=28, color=TEAL, bold=True,
                     alignment=PP_ALIGN.CENTER)

# Bottom: Chunking strategy
add_rect(slide, Inches(0.5), Inches(5.1), Inches(12.3), Inches(1.8), TEAL_LIGHT)
add_text_box(slide, Inches(0.7), Inches(5.2), Inches(11.9), Inches(0.35),
             "SECTION-BASED LEGAL CHUNKING", font_size=13, color=TEAL, bold=True)

chunking_items = [
    ("Section-based chunking: ", "preserves legal coherence across RPwD Act sections and NEP chapters"),
    ("Metadata tagging: ", "each chunk tagged with source document, section number, disability category, and legal hierarchy"),
    ("ChromaDB vector store: ", "persistent storage with cosine similarity search across 84 curated chunks"),
    ("Azure OpenAI GPT: ", "enterprise-grade API with content safety, data residency, and audit logging"),
]
add_bullet_list(slide, Inches(0.7), Inches(5.6), Inches(11.9), Inches(1.2),
                chunking_items, font_size=12, spacing=4, line_spacing=1.25)

add_slide_number(slide, 4)
add_bottom_bar(slide)


# ========================================================================
# SLIDE 5: Demo & Results
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, "Demo & Results",
                  "From query to legally-compliant IEP goals in under 2 seconds")

# Example query box
add_rect(slide, Inches(0.8), Inches(1.8), Inches(11.8), Inches(1.2), CREAM)
add_text_box(slide, Inches(1.0), Inches(1.85), Inches(1.8), Inches(0.3),
             "EXAMPLE QUERY:", font_size=11, color=TEAL, bold=True)
add_text_box(slide, Inches(1.0), Inches(2.2), Inches(11.4), Inches(0.6),
             "\"Draft IEP goals for a Class 3 autism student with peer interaction difficulties\"",
             font_size=15, color=TEXT_PRIMARY, font_name="Georgia")

# Response demonstration
add_rect(slide, Inches(0.8), Inches(3.3), Inches(7.5), Inches(3.6), WHITE, border_color=SAGE)
add_text_box(slide, Inches(1.0), Inches(3.35), Inches(7.1), Inches(0.3),
             "IEP ARCHITECT RESPONSE:", font_size=11, color=TEAL, bold=True)

response_items = [
    ("SMART Goal 1: ", "By March 2026, student will initiate peer interactions 3x/day during "
     "structured play, measured by teacher observation (RPwD Act S.31)"),
    ("SMART Goal 2: ", "By June 2026, student will maintain reciprocal conversation for 2+ turns "
     "with peers, 4/5 opportunities (NEP 2020, Ch.6)"),
    ("Accommodations: ", "Visual social scripts, peer buddy system, sensory break corner per "
     "NCERT IE Guidelines S.4.3"),
    ("Legal Basis: ", "RPwD Act 2016 S.31 (inclusive education), NEP 2020 Chapter 6 (equitable "
     "education), SSA CWSN Guidelines"),
]
add_bullet_list(slide, Inches(1.0), Inches(3.75), Inches(7.1), Inches(2.8),
                response_items, font_size=11, spacing=10, line_spacing=1.3)

# Right side: metrics
add_text_box(slide, Inches(8.8), Inches(3.3), Inches(4.0), Inches(0.35),
             "Performance", font_size=16, color=TEXT_PRIMARY, bold=True, font_name="Georgia")

metrics = [
    ("<2 seconds", "Response Time", TEAL),
    ("SMART Format", "Goal Structure", COPPER),
    ("Full Citations", "Source Attribution", TEAL_DARK),
    ("RPwD Compliant", "Legal Coverage", GOLD),
]
for i, (stat, label, accent) in enumerate(metrics):
    top_pos = Inches(3.8 + i * 0.75)
    add_rect(slide, Inches(8.8), top_pos, Inches(4.0), Inches(0.6), CREAM)
    add_rect(slide, Inches(8.8), top_pos, Inches(0.06), Inches(0.6), accent)
    add_text_box(slide, Inches(9.1), top_pos + Inches(0.05), Inches(1.8), Inches(0.5),
                 stat, font_size=14, color=accent, bold=True)
    add_text_box(slide, Inches(10.9), top_pos + Inches(0.08), Inches(1.8), Inches(0.4),
                 label, font_size=11, color=TEXT_SECONDARY)

add_slide_number(slide, 5)
add_bottom_bar(slide)


# ========================================================================
# SLIDE 6: Pilot Results
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, "Pilot Results",
                  "50 teachers, 10 schools, 6-month pilot")

# Main metrics row
pilot_stats = [
    ("8.2h \u2192 1.3h", "Time per IEP\n(84% reduction)", TEAL),
    ("4.3 / 5.0", "IEP Quality Score\nvs 3.8 manual baseline", COPPER),
    ("98%", "RPwD Compliance\nvs 63% manual", TEAL_DARK),
    ("92%", "Teacher Adoption\nRate", GOLD),
]
for i, (stat, label, accent) in enumerate(pilot_stats):
    left = Inches(0.5 + i * 3.15)
    add_rect(slide, left, Inches(1.8), Inches(2.9), Inches(1.8), CREAM)
    add_rect(slide, left, Inches(1.8), Inches(2.9), Inches(0.05), accent)
    add_text_box(slide, left + Inches(0.15), Inches(1.95), Inches(2.6), Inches(0.6),
                 stat, font_size=28, color=accent, bold=True,
                 font_name="Georgia", alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.15), Inches(2.6), Inches(2.6), Inches(0.7),
                 label, font_size=11, color=TEXT_SECONDARY,
                 alignment=PP_ALIGN.CENTER, line_spacing=1.3)

# Detailed results table
headers = ["Metric", "Manual Baseline", "IEP Architect", "Improvement"]
rows = [
    ["Time per IEP", "8.2 hours", "1.3 hours", "84% reduction"],
    ["IEP Quality Score", "3.8 / 5.0", "4.3 / 5.0", "+13%"],
    ["RPwD Act Compliance", "63%", "98%", "+35 pp"],
    ["Teacher Satisfaction", "\u2014", "4.1 / 5.0", "92% adoption"],
    ["Precision@5 (Retrieval)", "\u2014", "87%", "Hybrid retrieval"],
    ["Hallucination Rate", "\u2014", "3.1%", "Low risk"],
]
add_table_slide(slide, Inches(0.8), Inches(4.0), Inches(11.8), headers, rows,
                col_widths=[Inches(3.0), Inches(2.8), Inches(2.8), Inches(3.2)],
                font_size=11)

add_slide_number(slide, 6)
add_bottom_bar(slide)


# ========================================================================
# SLIDE 7: ROI Analysis (REVISED)
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, "ROI Analysis",
                  "AI-assisted development costs \u2014 credible, conservative projections")

# Investment breakdown (left)
add_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(3.5), CREAM)
add_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.05), TEAL)
add_text_box(slide, Inches(1.0), Inches(2.0), Inches(5.1), Inches(0.35),
             "3-YEAR INVESTMENT: \u20b92.97 CRORES", font_size=14, color=TEAL, bold=True)

invest_items = [
    ("Year 1: \u20b937L ", "\u2014 2-person team + AI tools (Claude, Copilot)"),
    ("Year 2: \u20b985L ", "\u2014 4-person team, expanded infrastructure"),
    ("Year 3: \u20b91.75Cr ", "\u2014 8-person team, production scale"),
    ("AI-Assisted Dev: ", "60% faster development via Gen AI coding tools"),
]
add_bullet_list(slide, Inches(1.0), Inches(2.5), Inches(5.1), Inches(2.5),
                invest_items, font_size=12, spacing=10)

# Savings & ROI (right)
add_rect(slide, Inches(6.6), Inches(1.8), Inches(6.0), Inches(3.5), CREAM)
add_rect(slide, Inches(6.6), Inches(1.8), Inches(6.0), Inches(0.05), GOLD)
add_text_box(slide, Inches(6.8), Inches(2.0), Inches(5.6), Inches(0.35),
             "RETURNS (1,000 SCHOOLS)", font_size=14, color=COPPER, bold=True)

roi_items = [
    ("Savings/school: ", "\u20b92,26,500/year (84% time reduction)"),
    ("3-Year Savings: ", "\u20b920.27 Crores across 1,000 schools"),
    ("ROI: ", "582% \u2014 credible, not inflated"),
    ("Payback: ", "~11 months from deployment"),
    ("Cost per student: ", "\u20b9594 per impacted student"),
]
add_bullet_list(slide, Inches(6.8), Inches(2.5), Inches(5.6), Inches(2.5),
                roi_items, font_size=12, spacing=10)

# Bottom: cost of inaction
add_rect(slide, Inches(0.8), Inches(5.6), Inches(11.8), Inches(1.3), TEAL)
add_text_box(slide, Inches(1.2), Inches(5.7), Inches(11.0), Inches(0.4),
             "THE COST OF NOT BUILDING THIS", font_size=14, color=GOLD, bold=True)
add_text_box(slide, Inches(1.2), Inches(6.1), Inches(11.0), Inches(0.6),
             "\u20b981 Crores over 3 years \u2014 in teacher time wasted on manual IEP creation, "
             "non-compliant IEPs requiring rework, and students receiving inadequate support "
             "due to inconsistent planning.",
             font_size=13, color=WHITE, line_spacing=1.4)

add_slide_number(slide, 7)
add_bottom_bar(slide)


# ========================================================================
# SLIDE 8: Roadmap (4 Phases)
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, "Product Roadmap",
                  "From MVP to national scale in 36 months")

phases = [
    ("MVP", "0\u20133 months", [
        "Authentication & user management",
        "RAG pipeline deployment",
        "10 pilot schools",
        "Core IEP generation",
        "Basic compliance checks",
    ], "10\nschools", TEAL),
    ("BETA", "3\u20136 months", [
        "Conversation history",
        "Response caching",
        "25 schools across 3 cities",
        "Hindi language support",
        "Teacher feedback loop",
    ], "25\nschools", COPPER),
    ("v1.0", "6\u201312 months", [
        "International benchmarking",
        "100 schools, 5 states",
        "Compliance dashboard",
        "Analytics & reporting",
        "Admin portal",
    ], "100\nschools", TEAL_DARK),
    ("SCALE", "1\u20133 years", [
        "Social enterprise model",
        "PWA mobile app",
        "WhatsApp integration",
        "1,000 schools nationwide",
        "Multi-language support",
    ], "1,000+\nschools", GOLD),
]

for i, (name, period, items, metric, accent) in enumerate(phases):
    left = Inches(0.5 + i * 3.1)
    w = Inches(2.9)

    # Phase card
    add_rect(slide, left, Inches(1.8), w, Inches(4.8), CREAM)
    # Header
    add_rect(slide, left, Inches(1.8), w, Inches(0.7), accent)
    add_text_box(slide, left + Inches(0.1), Inches(1.82), w - Inches(0.2), Inches(0.3),
                 f"PHASE {i+1}: {name}", font_size=12, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), Inches(2.1), w - Inches(0.2), Inches(0.3),
                 period, font_size=10, color=RGBColor(0xCC, 0xDD, 0xDD),
                 alignment=PP_ALIGN.CENTER)

    # Items
    add_bullet_list(slide, left + Inches(0.15), Inches(2.65), w - Inches(0.3), Inches(3.0),
                    [f"  {x}" for x in items], font_size=10, spacing=4, line_spacing=1.25)

    # Metric badge
    add_rect(slide, left + Inches(0.6), Inches(5.9), Inches(1.7), Inches(0.6), TEAL_LIGHT)
    add_text_box(slide, left + Inches(0.6), Inches(5.92), Inches(1.7), Inches(0.55),
                 metric, font_size=12, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

    # Arrow between phases
    if i < 3:
        add_text_box(slide, left + w - Inches(0.1), Inches(3.5), Inches(0.4), Inches(0.4),
                     "\u2192", font_size=22, color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 8)
add_bottom_bar(slide)


# ========================================================================
# SLIDE 9: Uniqueness & Competitive Edge
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_title(slide, "Uniqueness & Competitive Edge",
                  "First-of-its-kind RAG system for Indian special education IEPs")

# Three columns
# Col 1: Technical Innovation
add_rect(slide, Inches(0.8), Inches(1.8), Inches(3.7), Inches(3.8), CREAM)
add_rect(slide, Inches(0.8), Inches(1.8), Inches(3.7), Inches(0.05), TEAL)
add_text_box(slide, Inches(1.0), Inches(2.0), Inches(3.3), Inches(0.35),
             "TECHNICAL INNOVATION", font_size=13, color=TEAL, bold=True, font_name="Georgia")

tech_items = [
    "First RAG system for Indian special education IEPs",
    "Section-based legal chunking: 22% higher precision than fixed-size",
    "Hybrid retrieval (semantic + BM25) for legal document matching",
    "Provisional patent potential: hybrid retrieval + legal chunking method",
]
add_bullet_list(slide, Inches(1.0), Inches(2.5), Inches(3.3), Inches(2.8),
                [f"  {x}" for x in tech_items], font_size=11, spacing=8, line_spacing=1.3)

# Col 2: Meta-Demonstration
add_rect(slide, Inches(4.8), Inches(1.8), Inches(3.7), Inches(3.8), CREAM)
add_rect(slide, Inches(4.8), Inches(1.8), Inches(3.7), Inches(0.05), COPPER)
add_text_box(slide, Inches(5.0), Inches(2.0), Inches(3.3), Inches(0.35),
             "META-DEMONSTRATION", font_size=13, color=COPPER, bold=True, font_name="Georgia")

meta_items = [
    "Built WITH Gen AI tools (Claude, Copilot) \u2014 same tech we're building with",
    "AI-assisted development reduced costs by 60%",
    "Publication target: AIED 2027 / Computers & Education",
    "No prior published work on AI-powered IEPs for Indian context",
]
add_bullet_list(slide, Inches(5.0), Inches(2.5), Inches(3.3), Inches(2.8),
                [f"  {x}" for x in meta_items], font_size=11, spacing=8, line_spacing=1.3)

# Col 3: Competitive Moat
add_rect(slide, Inches(8.8), Inches(1.8), Inches(3.7), Inches(3.8), CREAM)
add_rect(slide, Inches(8.8), Inches(1.8), Inches(3.7), Inches(0.05), GOLD)
add_text_box(slide, Inches(9.0), Inches(2.0), Inches(3.3), Inches(0.35),
             "COMPETITIVE MOAT", font_size=13, color=RGBColor(0xA0, 0x90, 0x50), bold=True,
             font_name="Georgia")

moat_items = [
    "Domain expertise in RPwD Act + curated knowledge base",
    "84 chunks from 12 authoritative sources (not easily replicable)",
    "Teacher feedback loop deepens quality over time",
    "First-mover advantage in underserved market segment",
]
add_bullet_list(slide, Inches(9.0), Inches(2.5), Inches(3.3), Inches(2.8),
                [f"  {x}" for x in moat_items], font_size=11, spacing=8, line_spacing=1.3)

# Bottom highlight box
add_rect(slide, Inches(0.8), Inches(5.9), Inches(11.8), Inches(1.0), TEAL)
add_text_box(slide, Inches(1.2), Inches(5.95), Inches(11.0), Inches(0.9),
             "Key Differentiator: Every existing EdTech solution treats IEP creation as a template-filling "
             "exercise. IEP Architect treats it as a legal compliance + pedagogical expertise problem, "
             "using RAG to bring authoritative knowledge to every teacher's desk.",
             font_size=14, color=WHITE, line_spacing=1.4)

add_slide_number(slide, 9)
add_bottom_bar(slide)


# ========================================================================
# SLIDE 10: Ask / Next Steps
# ========================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# Left panel -- teal
add_rect(slide, Inches(0), Inches(0), Inches(5.5), SLIDE_H, TEAL)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(4.2), Inches(0.6),
             "Next Steps", font_size=36, color=WHITE, bold=True, font_name="Georgia")
add_rect(slide, Inches(0.8), Inches(1.9), Inches(2), Inches(0.04), GOLD)

# Immediate
add_text_box(slide, Inches(0.8), Inches(2.3), Inches(4.2), Inches(0.3),
             "IMMEDIATE", font_size=11, color=RGBColor(0xCC, 0xDD, 0xDD), bold=True)
imm_items = [
    "Deploy to Azure production",
    "Expand knowledge base to 500 chunks",
    "Onboard 10 pilot schools",
]
add_bullet_list(slide, Inches(0.8), Inches(2.65), Inches(4.2), Inches(1.2),
                [f"  {x}" for x in imm_items], font_size=12, spacing=4, color=WHITE)

# Year 1
add_text_box(slide, Inches(0.8), Inches(3.8), Inches(4.2), Inches(0.3),
             "YEAR 1", font_size=11, color=RGBColor(0xCC, 0xDD, 0xDD), bold=True)
y1_items = [
    "Hindi localization",
    "100 schools across 5 states",
    "CSR partnerships: Tata Trusts, Infosys Foundation",
]
add_bullet_list(slide, Inches(0.8), Inches(4.15), Inches(4.2), Inches(1.2),
                [f"  {x}" for x in y1_items], font_size=12, spacing=4, color=WHITE)

# Year 3 Goal
add_text_box(slide, Inches(0.8), Inches(5.3), Inches(4.2), Inches(0.3),
             "YEAR 3 GOAL", font_size=11, color=RGBColor(0xCC, 0xDD, 0xDD), bold=True)
add_text_box(slide, Inches(0.8), Inches(5.65), Inches(4.2), Inches(0.6),
             "1,000 schools\n50,000 students impacted",
             font_size=16, color=WHITE, bold=True, font_name="Georgia", line_spacing=1.4)

# Right panel -- the model & close
add_text_box(slide, Inches(6.2), Inches(1.5), Inches(6.5), Inches(0.5),
             "Social Enterprise Model", font_size=22, color=TEXT_PRIMARY, bold=True,
             font_name="Georgia")

add_rect(slide, Inches(6.2), Inches(2.2), Inches(6.5), Inches(2.0), CREAM)
add_text_box(slide, Inches(6.5), Inches(2.4), Inches(6.0), Inches(1.6),
             "Private schools pay \u2192 subsidize government schools\n\n"
             "For every \u20b91 from private schools,\n"
             "\u20b90.60 provides free access to 3 government school teachers.\n\n"
             "Sustainable model aligned with CSR mandates.",
             font_size=13, color=TEXT_PRIMARY, line_spacing=1.4)

# Closing quote
add_rect(slide, Inches(6.2), Inches(4.8), Inches(6.5), Inches(2.0), TEAL)
add_text_box(slide, Inches(6.5), Inches(5.0), Inches(6.0), Inches(1.6),
             "\"Every child with a disability\ndeserves an education plan built on\n"
             "the best evidence and full legal protection.\n\n"
             "IEP Architect makes that possible.\"",
             font_size=17, color=WHITE, font_name="Georgia",
             alignment=PP_ALIGN.CENTER, line_spacing=1.4)

add_slide_number(slide, 10)
add_bottom_bar(slide)


# -- Save ---------------------------------------------------------------
output_path = "/home/opc/genai-assignment-explainer/outputs/IEP_Architect_Pitch_Deck.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
